"""파일 형식별 텍스트 추출 서비스"""
import io
import logging

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
TEXT_EXTENSIONS = {".sql", ".json", ".txt", ".csv", ".md", ".html", ".htm"}


def is_binary_file(filename: str) -> bool:
    """바이너리 파싱이 필요한 파일인지 확인"""
    if not filename:
        return False
    ext = _get_ext(filename)
    return ext in SUPPORTED_EXTENSIONS


def _get_ext(filename: str) -> str:
    return "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def extract_text(filename: str, content_bytes: bytes) -> str:
    """파일에서 텍스트 추출"""
    ext = _get_ext(filename)

    if ext == ".pdf":
        return _extract_pdf(content_bytes)
    elif ext in (".docx", ".doc"):
        return _extract_docx(content_bytes)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(content_bytes)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(content_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")


def _extract_pdf(content_bytes: bytes) -> str:
    """PDF에서 텍스트 추출"""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content_bytes, filetype="pdf")
        texts = []
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                texts.append(f"[페이지 {page_num + 1}]\n{text.strip()}")
        doc.close()

        if not texts:
            return "(PDF에서 텍스트를 추출할 수 없습니다. 이미지 기반 PDF일 수 있습니다.)"

        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PDF 추출 실패: {e}")
        raise ValueError(f"PDF 파일을 읽을 수 없습니다: {e}")


def _extract_docx(content_bytes: bytes) -> str:
    """DOCX에서 텍스트 추출"""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content_bytes))
        texts = []

        # 본문 단락
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())

        # 표
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(row_text)
            if table_text:
                texts.append("\n".join(table_text))

        if not texts:
            return "(DOCX에서 텍스트를 추출할 수 없습니다.)"

        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"DOCX 추출 실패: {e}")
        raise ValueError(f"DOCX 파일을 읽을 수 없습니다: {e}")


def _extract_pptx(content_bytes: bytes) -> str:
    """PPTX에서 텍스트 추출"""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content_bytes))
        texts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        slide_texts.append(row_text)

            if slide_texts:
                texts.append(f"[슬라이드 {slide_num + 1}]\n" + "\n".join(slide_texts))

        if not texts:
            return "(PPTX에서 텍스트를 추출할 수 없습니다.)"

        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PPTX 추출 실패: {e}")
        raise ValueError(f"PPTX 파일을 읽을 수 없습니다: {e}")


def _extract_xlsx(content_bytes: bytes) -> str:
    """XLSX에서 텍스트 추출"""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
        texts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_texts = []

            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(v for v in row_values if v)
                if row_text.strip():
                    sheet_texts.append(row_text)

            if sheet_texts:
                texts.append(f"[시트: {sheet_name}]\n" + "\n".join(sheet_texts))

        wb.close()

        if not texts:
            return "(XLSX에서 데이터를 추출할 수 없습니다.)"

        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"XLSX 추출 실패: {e}")
        raise ValueError(f"XLSX 파일을 읽을 수 없습니다: {e}")
